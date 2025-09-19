import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import subprocess


def custom_image_sort(filename):
    """
    自定义图片排序函数：
    1. 普通图片 (fig_) 优先于补充图片 (SI_fig_)
    2. 按数字大小排序
    
    参数:
    filename: 图片文件名
    
    返回:
    排序用的元组 (类型, 数字)
    """
    # 尝试匹配 fig_数字.png 格式
    match = re.match(r'fig_(\d+)\.png$', filename, re.IGNORECASE)
    if match:
        return (0, int(match.group(1)))  # 类型0表示普通图片
    
    # 尝试匹配 SI_fig_数字.png 格式
    match = re.match(r'SI_fig_(\d+)\.png$', filename, re.IGNORECASE)
    if match:
        return (1, int(match.group(1)))  # 类型1表示补充图片
    
    # 其他格式的图片放在最后
    return (2, 0)


def get_sorted_images(image_folder):
    """
    获取排序后的图片文件列表
    
    参数:
    image_folder: 图片文件夹路径
    
    返回:
    排序后的图片文件列表
    """
    image_files = [f for f in os.listdir(image_folder) 
                  if f.lower().endswith('.png')]
    image_files.sort(key=custom_image_sort)
    return image_files


def create_slide_with_image(presentation, image_path, description=None):
    """
    创建包含图片和描述的幻灯片
    
    参数:
    presentation: PPT对象
    image_path: 图片完整路径
    description: 图片描述文本(可选）
    """
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    
    # 获取幻灯片尺寸
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    
    # 获取图片原始尺寸
    with Image.open(image_path) as img:
        img_width, img_height = img.size
        img_aspect_ratio = img_width / img_height
    
    # 计算最大允许尺寸(80%的幻灯片尺寸）
    max_width = slide_width * 0.8
    max_height = slide_height * 0.8
    
    # 计算保持宽高比的目标尺寸
    target_width = min(max_width, img_aspect_ratio * max_height)
    target_height = min(max_height, max_width / img_aspect_ratio)
    
    # 计算居中位置
    left = (slide_width - target_width) / 2
    top = (slide_height - target_height) / 2
    
    # 添加描述文本框(如果有描述）在图片上方
    if description:
        # 文本框位置(图片上方）
        textbox_height = Inches(1.5)  # 增加文本框高度以适应大字体
        textbox_top = top - textbox_height - Inches(0.2)  # 在图片上方0.2英寸处
        
        # 确保文本框不会超出幻灯片顶部
        if textbox_top < Inches(0.1):
            textbox_top = Inches(0.1)
        
        # 创建文本框
        textbox = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=textbox_top,
            width=slide_width - Inches(1),
            height=textbox_height
        )
        
        # 设置文本框格式
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # 添加段落并设置文本
        p = text_frame.add_paragraph()
        p.text = description
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(30)  # 修改为30号字体
        p.font.name = 'Arial'
        p.font.bold = False
        p.font.color.rgb = RGBColor(0, 0, 0)  # 黑色文本
    
    # 添加图片到幻灯片
    slide.shapes.add_picture(
        image_path, 
        left, 
        top,
        width=target_width,
        height=target_height
    )

    return slide


def create_presentation_from_images(image_folder, output_filename, image_descriptions=None):
    """
    从指定文件夹读取按特定规则命名的图片生成PPT
    
    参数:
    image_folder: 图片文件夹路径
    output_filename: 输出的PPT文件名
    image_descriptions: 图片描述字典(key: 文件名, value: 描述）
    """
    prs = Presentation()
    
    # 设置幻灯片比例为16:9
    prs.slide_width = Inches(13.333)  # 16:9宽度
    prs.slide_height = Inches(7.5)    # 16:9高度

    image_files = get_sorted_images(image_folder)
    
    for filename in image_files:
        img_path = os.path.join(image_folder, filename)
        
        # 获取当前图片的描述(如果提供了描述字典）
        description = None
        if image_descriptions:
            # 使用文件名作为键获取描述,如果找不到则尝试小写匹配
            description = image_descriptions.get(filename) or \
                          image_descriptions.get(filename.lower())
        
        create_slide_with_image(prs, img_path, description)
    
    prs.save(output_filename)
    print(f"成功创建PPT: {output_filename},包含 {len(image_files)} 张图片")
    convert_pptx_to_pdf(output_filename, os.path.dirname(output_filename))


def convert_pptx_to_pdf(input_path, output_dir=None):
    if output_dir is None:
        output_dir = os.path.dirname(input_path)
    
    cmd = [
        'soffice', 
        '--headless', 
        '--convert-to', 
        'pdf', 
        '--outdir', 
        output_dir, 
        input_path
    ]
    
    subprocess.run(cmd, check=True)
    return os.path.join(output_dir, os.path.splitext(os.path.basename(input_path))[0] + ".pdf")